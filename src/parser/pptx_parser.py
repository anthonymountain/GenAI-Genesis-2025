import os
from pptx import Presentation
from pptx.util import Inches


def create_pptx_from_response(response_text, output_path="presentation.pptx",
                              slide_images=None):
    """
    Converts structured LLM output into a PowerPoint presentation.
    Optionally inserts one image per slide.
    """
    prs = Presentation()
    slides = response_text.split("Slide")[1:]
    slide_images = slide_images or []

    for i, raw_slide in enumerate(slides):
        lines = raw_slide.strip().split("\n")
        title = ""
        bullets = []
        script = ""
        reading_script = False

        for line in lines:

            if ":" in line and "title" in line.lower().split(":")[0]:
                title = line.split(":", 1)[-1].strip()
            elif line.startswith("-"):
                bullets.append(line.strip())
            elif line.startswith("Presenter Script:"):
                script = line.replace("Presenter Script:", "").strip()
                reading_script = True
            elif reading_script:
                script += " " + line.strip()

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        content = slide.placeholders[1]
        content.text = "\n".join(bullets)

        # Add speaker notes
        if script:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = script.strip()

        # Add image if available
        if i < len(slide_images) and slide_images[i]:
            image_path = slide_images[i]
            if image_path and os.path.exists(image_path):
                try:
                    with open(image_path, "rb") as img_file:
                        img_width = Inches(3)
                        img_height = Inches(2)
                        left = prs.slide_width - img_width - Inches(0.5)
                        top = prs.slide_height - img_height - Inches(0.5)
                        slide.shapes.add_picture(img_file, left, top,
                                                 width=img_width,
                                                 height=img_height)
                except Exception as e:
                    print(f"Failed to add image to slide {i+1}: {e}")

    prs.save(output_path)
    return output_path
